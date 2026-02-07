"""
Backend SaaS FastAPI - Camada de Plataforma
Responsável por: Autenticação, Multi-tenant, Billing, API
NÃO contém lógica de decisão (isso fica no core)
"""

from fastapi import FastAPI, Depends, HTTPException, status, BackgroundTasks, File, UploadFile, Form
from fastapi.middleware.cors import CORSMiddleware
from fastapi.security import HTTPBearer, HTTPAuthorizationCredentials
from typing import List
import asyncio
from pydantic import BaseModel, EmailStr
from datetime import datetime, timedelta
from typing import Optional
import jwt
import os
from pathlib import Path
from dotenv import load_dotenv

# Carrega .env do diretório atual e do diretório pai (raiz do projeto)
load_dotenv()  # Tenta backend/.env
load_dotenv(Path(__file__).parent.parent / ".env")  # Tenta raiz/.env

# ============================================================================
# CONFIGURAÇÃO
# ============================================================================

app = FastAPI(
    title="Consultor Executivo Multi-Agentes - API",
    version="1.0.0",
    description="API SaaS para análise estratégica com múltiplos agentes"
)

# CORS
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # Em produção: ["https://app.example.com"]
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Configurações
SECRET_KEY = os.getenv("JWT_SECRET_KEY", "dev-secret-key-change-in-production")
ALGORITHM = "HS256"
ACCESS_TOKEN_EXPIRE_MINUTES = 15
REFRESH_TOKEN_EXPIRE_DAYS = 30

security = HTTPBearer()

# ============================================================================
# MODELOS
# ============================================================================

class UserRegister(BaseModel):
    email: EmailStr
    password: str
    organization_name: str


class UserLogin(BaseModel):
    email: EmailStr
    password: str


class TokenResponse(BaseModel):
    access_token: str
    refresh_token: str
    token_type: str = "bearer"


class OrganizationInfo(BaseModel):
    id: str
    name: str
    plan: str

class UserResponse(BaseModel):
    id: Optional[str] = None
    user_id: str
    email: str
    name: Optional[str] = None
    role: str = "user"
    org_id: Optional[str] = None
    tenant_id: str
    organization: Optional[OrganizationInfo] = None
    organization_name: str
    created_at: datetime


class ExecutionRequest(BaseModel):
    problem_description: str
    business_type: str = "B2B"
    analysis_depth: str = "Padrão"


class ExecutionResponse(BaseModel):
    execution_id: str
    status: str
    created_at: datetime
    estimated_duration_seconds: int = 30


class BillingStatus(BaseModel):
    plan: str  # free, pro, enterprise
    executions_used: int
    executions_limit: Optional[int]
    tokens_used_today: int
    tokens_limit_today: Optional[int]
    billing_status: str  # active, past_due, cancelled
    next_billing_date: Optional[datetime]


# ============================================================================
# AUTENTICAÇÃO
# ============================================================================

def create_access_token(user_id: str, email: str, tenant_id: str) -> str:
    """Cria JWT access token"""
    payload = {
        "sub": user_id,
        "email": email,
        "tenant_id": tenant_id,
        "exp": datetime.utcnow() + timedelta(minutes=ACCESS_TOKEN_EXPIRE_MINUTES),
        "iat": datetime.utcnow(),
        "type": "access"
    }
    return jwt.encode(payload, SECRET_KEY, algorithm=ALGORITHM)


def create_refresh_token(user_id: str) -> str:
    """Cria JWT refresh token"""
    payload = {
        "sub": user_id,
        "exp": datetime.utcnow() + timedelta(days=REFRESH_TOKEN_EXPIRE_DAYS),
        "iat": datetime.utcnow(),
        "type": "refresh"
    }
    return jwt.encode(payload, SECRET_KEY, algorithm=ALGORITHM)


def verify_token(token: str) -> dict:
    """Valida e decodifica JWT"""
    try:
        payload = jwt.decode(token, SECRET_KEY, algorithms=[ALGORITHM])
        return payload
    except jwt.ExpiredSignatureError:
        raise HTTPException(
            status_code=status.HTTP_401_UNAUTHORIZED,
            detail="Token expirado"
        )
    except jwt.InvalidTokenError:
        raise HTTPException(
            status_code=status.HTTP_401_UNAUTHORIZED,
            detail="Token inválido"
        )


async def get_current_user(credentials: HTTPAuthorizationCredentials = Depends(security)) -> dict:
    """Middleware de autenticação"""
    token = credentials.credentials
    payload = verify_token(token)
    
    if payload.get("type") != "access":
        raise HTTPException(
            status_code=status.HTTP_401_UNAUTHORIZED,
            detail="Token inválido"
        )
    
    return payload


# ============================================================================
# TENANT CONTEXT
# ============================================================================

class TenantContext:
    """Contexto do tenant extraído do token"""
    def __init__(self, tenant_id: str, user_id: str, email: str):
        self.tenant_id = tenant_id
        self.user_id = user_id
        self.email = email


async def get_tenant_context(current_user: dict = Depends(get_current_user)) -> TenantContext:
    """Extrai contexto do tenant do token"""
    return TenantContext(
        tenant_id=current_user["tenant_id"],
        user_id=current_user["sub"],
        email=current_user["email"]
    )


# ============================================================================
# BILLING CONTROL
# ============================================================================

class BillingService:
    """Mock de serviço de billing - Em produção seria real"""
    
    # Limites por plano (alinhados com a UI)
    PLAN_LIMITS = {
        "free": {
            "analyses_per_month": 5,
            "agents": 3,  # analyst, commercial, reviewer
            "export_formats": [],  # Não pode exportar
            "refine_messages_per_analysis": 3,  # Máximo 3 perguntas por análise
            "features": ["Suporte por email", "3 perguntas de refino por análise"]
        },
        "pro": {
            "analyses_per_month": 50,
            "agents": 5,  # Todos
            "export_formats": ["pdf"],  # Só PDF
            "refine_messages_per_analysis": 20,  # Até 20 perguntas por análise
            "features": ["Suporte prioritário", "Exportação PDF", "20 perguntas de refino por análise"]
        },
        "enterprise": {
            "analyses_per_month": -1,  # Ilimitado
            "agents": 5,
            "export_formats": ["pdf", "docx", "pptx"],  # Todos
            "refine_messages_per_analysis": -1,  # Ilimitado
            "features": ["Suporte dedicado", "API access", "White-label", "Refino ilimitado"]
        }
    }
    
    # Arquivo para persistência dos dados de billing
    BILLING_FILE = Path(__file__).parent / "billing_data.json"
    
    # Cache em memória (carregado do arquivo)
    tenants = {}
    
    @classmethod
    def _load_from_file(cls):
        """Carrega dados de billing do arquivo JSON"""
        if cls.BILLING_FILE.exists():
            try:
                import json
                with open(cls.BILLING_FILE, 'r') as f:
                    cls.tenants = json.load(f)
                print(f"📊 Billing data carregado: {len(cls.tenants)} tenants")
            except Exception as e:
                print(f"⚠️ Erro ao carregar billing data: {e}")
                cls.tenants = {}
        else:
            cls.tenants = {}
    
    @classmethod
    def _save_to_file(cls):
        """Salva dados de billing no arquivo JSON"""
        try:
            import json
            with open(cls.BILLING_FILE, 'w') as f:
                json.dump(cls.tenants, f, indent=2)
        except Exception as e:
            print(f"⚠️ Erro ao salvar billing data: {e}")
    
    @staticmethod
    def check_execution_allowed(tenant_id: str) -> tuple[bool, Optional[str]]:
        """Verifica se tenant pode executar análise"""
        tenant = BillingService.tenants.get(tenant_id)
        
        if not tenant:
            # Novo tenant: começa em free
            tenant = {
                "plan": "free",
                "executions_this_month": 0,
                "tokens_today": 0,
                "billing_status": "active"
            }
            BillingService.tenants[tenant_id] = tenant
            BillingService._save_to_file()  # Persiste novo tenant
        
        # Verifica status
        if tenant["billing_status"] != "active":
            return False, f"Assinatura {tenant['billing_status']}"
        
        plan = tenant["plan"]
        limits = BillingService.PLAN_LIMITS.get(plan, BillingService.PLAN_LIMITS["free"])
        
        # Verifica limite de análises por mês
        max_analyses = limits["analyses_per_month"]
        if max_analyses != -1:  # -1 = ilimitado
            if tenant["executions_this_month"] >= max_analyses:
                if plan == "free":
                    return False, f"Limite de {max_analyses} análises/mês atingido. Faça upgrade para Pro!"
                else:
                    return False, f"Limite de {max_analyses} análises/mês atingido. Faça upgrade para Enterprise!"
        
        return True, None
    
    @staticmethod
    def get_allowed_agents(tenant_id: str) -> list[str]:
        """Retorna lista de agentes permitidos para o plano"""
        tenant = BillingService.tenants.get(tenant_id, {"plan": "free"})
        plan = tenant.get("plan", "free")
        limits = BillingService.PLAN_LIMITS.get(plan, BillingService.PLAN_LIMITS["free"])
        
        if limits["agents"] == 3:
            # Free: apenas 3 agentes essenciais
            return ["analyst", "commercial", "reviewer"]
        else:
            # Pro/Enterprise: todos os 5 agentes
            return ["analyst", "commercial", "financial", "market", "reviewer"]
    
    @staticmethod
    def can_export(tenant_id: str, format: str) -> tuple[bool, Optional[str]]:
        """Verifica se tenant pode exportar no formato especificado"""
        tenant = BillingService.tenants.get(tenant_id, {"plan": "free"})
        plan = tenant.get("plan", "free")
        limits = BillingService.PLAN_LIMITS.get(plan, BillingService.PLAN_LIMITS["free"])
        
        allowed_formats = limits["export_formats"]
        
        if not allowed_formats:
            return False, "Exportação disponível apenas para planos Pro e Enterprise"
        
        if format not in allowed_formats:
            if plan == "pro":
                return False, f"Formato {format.upper()} disponível apenas no plano Enterprise"
            return False, f"Formato {format.upper()} não disponível para seu plano"
        
        return True, None
    
    @staticmethod
    def can_refine(tenant_id: str, analysis: dict) -> tuple[bool, int, int, Optional[str]]:
        """
        Verifica se tenant pode enviar mais mensagens de refino para uma análise.
        
        Returns:
            tuple: (pode_refinar, mensagens_usadas, limite, mensagem_erro)
        """
        tenant = BillingService.tenants.get(tenant_id, {"plan": "free"})
        plan = tenant.get("plan", "free")
        limits = BillingService.PLAN_LIMITS.get(plan, BillingService.PLAN_LIMITS["free"])
        
        max_messages = limits.get("refine_messages_per_analysis", 3)
        
        # Conta apenas mensagens do usuário no histórico
        chat_history = analysis.get("chat_history", [])
        user_messages = len([m for m in chat_history if m.get("role") == "user"])
        
        # -1 = ilimitado
        if max_messages == -1:
            return True, user_messages, -1, None
        
        if user_messages >= max_messages:
            if plan == "free":
                return False, user_messages, max_messages, f"Limite de {max_messages} perguntas atingido. Faça upgrade para Pro e tenha até 20 perguntas por análise!"
            elif plan == "pro":
                return False, user_messages, max_messages, f"Limite de {max_messages} perguntas atingido. Faça upgrade para Enterprise e tenha perguntas ilimitadas!"
            else:
                return False, user_messages, max_messages, f"Limite de {max_messages} perguntas atingido."
        
        return True, user_messages, max_messages, None
    
    @staticmethod
    def record_execution(tenant_id: str, tokens_used: int) -> None:
        """Registra execução para billing"""
        tenant = BillingService.tenants.get(tenant_id)
        if tenant:
            tenant["executions_this_month"] += 1
            tenant["tokens_today"] += tokens_used
            BillingService._save_to_file()  # Persiste após cada execução
    
    @staticmethod
    def get_status(tenant_id: str) -> BillingStatus:
        """Retorna status de billing do tenant"""
        tenant = BillingService.tenants.get(tenant_id, {"plan": "free"})
        plan = tenant.get("plan", "free")
        limits = BillingService.PLAN_LIMITS.get(plan, BillingService.PLAN_LIMITS["free"])
        
        return BillingStatus(
            plan=plan,
            executions_used=tenant["executions_this_month"],
            executions_limit=limits["analyses_per_month"] if limits["analyses_per_month"] != -1 else None,
            tokens_used_today=tenant.get("tokens_today", 0),
            tokens_limit_today=None,
            billing_status=tenant["billing_status"],
            next_billing_date=datetime.utcnow() + timedelta(days=30)
        )


# Carrega dados de billing persistidos ao iniciar
BillingService._load_from_file()


# ============================================================================
# ENDPOINTS: AUTENTICAÇÃO
# ============================================================================

@app.post("/api/v1/auth/register", response_model=TokenResponse)
async def register(user: UserRegister):
    """Registra novo usuário e organização"""
    # Em produção: validar email único, hash password, salvar BD
    user_id = f"user_{hash(user.email) % 10000}"
    tenant_id = f"org_{hash(user.organization_name) % 10000}"
    
    access_token = create_access_token(user_id, user.email, tenant_id)
    refresh_token = create_refresh_token(user_id)
    
    return TokenResponse(
        access_token=access_token,
        refresh_token=refresh_token
    )


@app.post("/api/v1/auth/login", response_model=TokenResponse)
async def login(user: UserLogin):
    """Login de usuário existente"""
    # Em produção: validar credenciais contra BD
    user_id = f"user_{hash(user.email) % 10000}"
    tenant_id = f"org_{hash(user.email) % 10000}"
    
    access_token = create_access_token(user_id, user.email, tenant_id)
    refresh_token = create_refresh_token(user_id)
    
    return TokenResponse(
        access_token=access_token,
        refresh_token=refresh_token
    )


@app.post("/api/v1/auth/refresh", response_model=TokenResponse)
async def refresh(refresh_token: str):
    """Renova access token usando refresh token"""
    payload = verify_token(refresh_token)
    
    if payload.get("type") != "refresh":
        raise HTTPException(
            status_code=status.HTTP_401_UNAUTHORIZED,
            detail="Refresh token inválido"
        )
    
    user_id = payload["sub"]
    # Em produção: recuperar email e tenant_id do BD
    email = "user@example.com"
    tenant_id = "org_123"
    
    access_token = create_access_token(user_id, email, tenant_id)
    
    return TokenResponse(
        access_token=access_token,
        refresh_token=refresh_token
    )


# ============================================================================
# ENDPOINTS: USUÁRIO
# ============================================================================

@app.get("/api/v1/me", response_model=UserResponse)
@app.get("/api/v1/auth/me", response_model=UserResponse)
async def get_current_user_info(tenant: TenantContext = Depends(get_tenant_context)):
    """Retorna informações do usuário atual"""
    return UserResponse(
        id=tenant.user_id,
        user_id=tenant.user_id,
        email=tenant.email,
        name=tenant.email.split("@")[0],
        role="admin",
        org_id=tenant.tenant_id,
        tenant_id=tenant.tenant_id,
        organization={
            "id": tenant.tenant_id,
            "name": "Minha Empresa",
            "plan": "free"
        },
        organization_name="Minha Empresa",
        created_at=datetime.utcnow()
    )


# ============================================================================
# ENDPOINTS: EXECUÇÕES
# ============================================================================

@app.post("/api/v1/executions", response_model=ExecutionResponse, status_code=201)
async def create_execution(
    request: ExecutionRequest,
    tenant: TenantContext = Depends(get_tenant_context)
):
    """Cria nova execução (análise)"""
    
    # 1. Verifica billing
    allowed, error_msg = BillingService.check_execution_allowed(tenant.tenant_id)
    if not allowed:
        raise HTTPException(
            status_code=status.HTTP_402_PAYMENT_REQUIRED,
            detail=error_msg
        )
    
    # 2. Chama core engine
    # Em produção: seria uma chamada real ao BusinessOrchestrator
    execution_id = f"exec_{hash(request.problem_description) % 100000}"
    
    # 3. Registra para billing
    BillingService.record_execution(tenant.tenant_id, tokens_used=500)
    
    return ExecutionResponse(
        execution_id=execution_id,
        status="running",
        created_at=datetime.utcnow(),
        estimated_duration_seconds=30
    )


@app.get("/api/v1/executions")
async def list_executions(tenant: TenantContext = Depends(get_tenant_context)):
    """Lista execuções do tenant"""
    # Em produção: query BD com WHERE tenant_id = ?
    return {
        "executions": [
            {
                "execution_id": "exec_123",
                "status": "completed",
                "created_at": datetime.utcnow().isoformat(),
                "problem": "Vendas caíram 20%..."
            }
        ]
    }


@app.get("/api/v1/executions/{execution_id}")
async def get_execution(
    execution_id: str,
    tenant: TenantContext = Depends(get_tenant_context)
):
    """Retorna detalhe de uma execução"""
    # Em produção: query BD com WHERE execution_id = ? AND tenant_id = ?
    return {
        "execution_id": execution_id,
        "status": "completed",
        "decision": "Investir em marketing digital",
        "confidence": 0.82,
        "actions": [
            {"description": "Preparar plano", "owner": "Commercial", "due": "5 dias"}
        ]
    }


# ============================================================================
# ENDPOINTS: ASYNC ANALYSES (para frontend)
# ============================================================================

# Armazenamento em memória para análises (em produção: usar BD)
analyses_store: dict = {}


def run_analysis_task(
    analysis_id: str,
    problem_description: str,
    business_type: str,
    analysis_depth: str,
    tenant_id: str
):
    """
    Executa análise com os agentes de IA em background.
    Atualiza o status da análise quando terminar.
    """
    try:
        # Atualiza status para running
        if analysis_id in analyses_store:
            analyses_store[analysis_id]["status"] = "running"
        
        # Importa o time de agentes
        from core.types import ExecutionContext
        from orchestrator import BusinessOrchestrator
        from agents import (
            AnalystAgent,
            CommercialAgent,
            FinancialAgent,
            MarketAgent,
            ReviewerAgent,
        )
        
        # Obtém agentes permitidos pelo plano
        allowed_agents = BillingService.get_allowed_agents(tenant_id)
        
        # Cria apenas os agentes permitidos pelo plano
        # ReviewerAgent recebe a lista de agentes disponíveis para ajustar dependências
        all_agents = {
            "analyst": AnalystAgent(),
            "commercial": CommercialAgent(),
            "financial": FinancialAgent(),
            "market": MarketAgent(),
            "reviewer": ReviewerAgent(available_agents=allowed_agents),
        }
        
        # Filtra agentes pelo plano (Free = 3, Pro/Enterprise = 5)
        agents = {k: v for k, v in all_agents.items() if k in allowed_agents}
        orchestrator = BusinessOrchestrator(agents)
        
        # Cria contexto
        context = ExecutionContext(
            problem_description=problem_description,
            business_type=business_type,
            analysis_depth=analysis_depth,
        )
        
        # Executa análise (converte async para sync)
        result_context = asyncio.run(orchestrator.execute(context))
        
        # Prepara resultados
        results = {
            "analyst": result_context.get_agent_output("analyst") or "",
            "commercial": result_context.get_agent_output("commercial") or "",
            "financial": result_context.get_agent_output("financial") or "",
            "market": result_context.get_agent_output("market") or "",
            "reviewer": result_context.get_agent_output("reviewer") or "",
        }
        
        # Atualiza análise com resultados
        if analysis_id in analyses_store:
            analyses_store[analysis_id].update({
                "status": "completed",
                "completed_at": datetime.utcnow().isoformat(),
                "executive_summary": results.get("reviewer", ""),
                "results": results,
                "agent_outputs": {
                    name: {"content": output, "output": output}
                    for name, output in results.items()
                },
                "total_tokens": result_context.get_total_tokens(),
                "total_cost_usd": result_context.get_total_cost(),
            })
        
        # Registra billing
        BillingService.record_execution(tenant_id, tokens_used=result_context.get_total_tokens())
        
        print(f"✅ Análise {analysis_id} concluída com sucesso!")
        
    except Exception as e:
        print(f"❌ Erro na análise {analysis_id}: {str(e)}")
        import traceback
        traceback.print_exc()
        
        # Marca como falha
        if analysis_id in analyses_store:
            analyses_store[analysis_id].update({
                "status": "failed",
                "error_message": str(e),
                "completed_at": datetime.utcnow().isoformat(),
            })

@app.get("/api/v1/async/analyses")
async def list_analyses(
    limit: int = 10,
    offset: int = 0,
    tenant: TenantContext = Depends(get_tenant_context)
):
    """Lista análises do tenant"""
    tenant_analyses = [
        a for a in analyses_store.values() 
        if a.get("tenant_id") == tenant.tenant_id
    ]
    # Ordenar por data de criação (mais recentes primeiro)
    tenant_analyses.sort(key=lambda x: x.get("created_at", ""), reverse=True)
    
    return {
        "items": tenant_analyses[offset:offset + limit],
        "total": len(tenant_analyses),
        "limit": limit,
        "offset": offset
    }


@app.post("/api/v1/async/analyses", status_code=201)
async def create_analysis(
    request: ExecutionRequest,
    background_tasks: BackgroundTasks,
    tenant: TenantContext = Depends(get_tenant_context)
):
    """Cria nova análise e executa em background"""
    # Verifica billing
    allowed, error_msg = BillingService.check_execution_allowed(tenant.tenant_id)
    if not allowed:
        raise HTTPException(
            status_code=status.HTTP_402_PAYMENT_REQUIRED,
            detail=error_msg
        )
    
    import random
    analysis_id = f"analysis_{random.randint(10000, 99999)}"
    created_at = datetime.utcnow().isoformat()
    
    # Criar análise com status PENDING
    analysis = {
        "id": analysis_id,
        "tenant_id": tenant.tenant_id,
        "problem_description": request.problem_description,
        "business_type": request.business_type,
        "analysis_depth": request.analysis_depth,
        "status": "pending",  # Começa como pendente
        "created_at": created_at,
        "executive_summary": None,
        "agent_outputs": {},
        "results": {}
    }
    
    # Armazenar análise
    analyses_store[analysis_id] = analysis
    
    # Executar análise em background
    background_tasks.add_task(
        run_analysis_task,
        analysis_id,
        request.problem_description,
        request.business_type,
        request.analysis_depth,
        tenant.tenant_id
    )
    
    return {"analysis_id": analysis_id, "status": "pending", "message": "Análise iniciada"}


@app.post("/api/v1/async/analyses/with-files", status_code=201)
async def create_analysis_with_files(
    background_tasks: BackgroundTasks,
    problem_description: str = Form(...),
    business_type: str = Form("SaaS"),
    analysis_depth: str = Form("Padrão"),
    files: List[UploadFile] = File(default=[]),
    tenant: TenantContext = Depends(get_tenant_context)
):
    """Cria nova análise com arquivos anexados"""
    # Verifica billing
    allowed, error_msg = BillingService.check_execution_allowed(tenant.tenant_id)
    if not allowed:
        raise HTTPException(
            status_code=status.HTTP_402_PAYMENT_REQUIRED,
            detail=error_msg
        )
    
    import random
    analysis_id = f"analysis_{random.randint(10000, 99999)}"
    created_at = datetime.utcnow().isoformat()
    
    # Processa arquivos e extrai conteúdo
    file_contents = []
    file_names = []
    
    for file in files:
        try:
            content = await file.read()
            file_names.append(file.filename)
            
            # Extrai texto baseado no tipo de arquivo
            if file.filename.endswith('.csv'):
                text = content.decode('utf-8', errors='ignore')
                file_contents.append(f"--- Dados do arquivo {file.filename} ---\n{text[:5000]}")
            
            elif file.filename.endswith('.txt'):
                text = content.decode('utf-8', errors='ignore')
                file_contents.append(f"--- Conteúdo de {file.filename} ---\n{text[:5000]}")
            
            elif file.filename.endswith(('.xlsx', '.xls')):
                # Extrai dados de planilhas Excel
                try:
                    import io
                    import openpyxl
                    
                    workbook = openpyxl.load_workbook(io.BytesIO(content), data_only=True)
                    excel_text = f"--- Planilha: {file.filename} ---\n"
                    
                    for sheet_name in workbook.sheetnames[:3]:  # Máximo 3 abas
                        sheet = workbook[sheet_name]
                        excel_text += f"\n📊 Aba: {sheet_name}\n"
                        
                        rows_data = []
                        for row_idx, row in enumerate(sheet.iter_rows(values_only=True)):
                            if row_idx > 50:  # Limita a 50 linhas
                                excel_text += f"... (mais {sheet.max_row - 50} linhas)\n"
                                break
                            # Converte valores para string
                            row_str = " | ".join([str(cell) if cell is not None else "" for cell in row[:10]])  # Máx 10 colunas
                            rows_data.append(row_str)
                        
                        excel_text += "\n".join(rows_data)
                    
                    file_contents.append(excel_text[:5000])
                except ImportError:
                    file_contents.append(f"--- Planilha {file.filename}: openpyxl não instalado ---")
                except Exception as e:
                    file_contents.append(f"--- Erro ao processar Excel {file.filename}: {str(e)} ---")
            
            elif file.filename.endswith('.pdf'):
                # Extrai texto de PDFs
                try:
                    import io
                    from PyPDF2 import PdfReader
                    
                    pdf_reader = PdfReader(io.BytesIO(content))
                    pdf_text = f"--- PDF: {file.filename} ({len(pdf_reader.pages)} páginas) ---\n"
                    
                    for page_num, page in enumerate(pdf_reader.pages[:10]):  # Máximo 10 páginas
                        page_text = page.extract_text()
                        if page_text:
                            pdf_text += f"\n📄 Página {page_num + 1}:\n{page_text[:1000]}\n"
                    
                    file_contents.append(pdf_text[:5000])
                except ImportError:
                    file_contents.append(f"--- PDF {file.filename}: PyPDF2 não instalado ---")
                except Exception as e:
                    file_contents.append(f"--- Erro ao processar PDF {file.filename}: {str(e)} ---")
            
            else:
                file_contents.append(f"--- Arquivo anexado: {file.filename} ---")
        except Exception as e:
            file_contents.append(f"--- Erro ao processar {file.filename}: {str(e)} ---")
    
    # Enriquece a descrição do problema com conteúdo dos arquivos
    enriched_description = problem_description
    if file_contents:
        enriched_description += "\n\n=== DADOS ANEXADOS ===\n" + "\n\n".join(file_contents)
    
    # Criar análise com status PENDING
    analysis = {
        "id": analysis_id,
        "tenant_id": tenant.tenant_id,
        "problem_description": problem_description,
        "enriched_description": enriched_description,
        "business_type": business_type,
        "analysis_depth": analysis_depth,
        "status": "pending",
        "created_at": created_at,
        "files": file_names,
        "executive_summary": None,
        "agent_outputs": {},
        "results": {}
    }
    
    # Armazenar análise
    analyses_store[analysis_id] = analysis
    
    # Executar análise em background (usa descrição enriquecida)
    background_tasks.add_task(
        run_analysis_task,
        analysis_id,
        enriched_description,  # Usa descrição enriquecida com dados dos arquivos
        business_type,
        analysis_depth,
        tenant.tenant_id
    )
    
    return {
        "analysis_id": analysis_id, 
        "status": "pending", 
        "message": "Análise iniciada",
        "files_processed": len(file_names)
    }


class RefineRequest(BaseModel):
    analysis_id: str
    message: str
    context: dict = {}


@app.post("/api/v1/async/analyses/refine")
async def refine_analysis(
    request: RefineRequest,
    tenant: TenantContext = Depends(get_tenant_context)
):
    """Refina análise com perguntas de follow-up"""
    analysis = analyses_store.get(request.analysis_id)
    
    if not analysis:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="Análise não encontrada"
        )
    
    if analysis.get("tenant_id") != tenant.tenant_id:
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="Acesso negado"
        )
    
    # Verifica limite de mensagens de refino pelo plano
    can_refine, used, limit, error_msg = BillingService.can_refine(tenant.tenant_id, analysis)
    if not can_refine:
        raise HTTPException(
            status_code=status.HTTP_402_PAYMENT_REQUIRED,
            detail={
                "message": error_msg,
                "used": used,
                "limit": limit,
                "upgrade_url": "/billing"
            }
        )
    
    try:
        import anthropic
        
        api_key = os.getenv("ANTHROPIC_API_KEY")
        if not api_key:
            return {"response": "API key não configurada. Configure ANTHROPIC_API_KEY no arquivo .env"}
        
        client = anthropic.Anthropic(api_key=api_key)
        
        # Monta contexto da análise anterior
        context_text = ""
        if analysis.get("results"):
            for agent, content in analysis["results"].items():
                if content:
                    context_text += f"\n\n### Análise do {agent}:\n{content[:2000]}"
        
        system_prompt = f"""Você é um consultor executivo experiente que está refinando uma análise de negócios.

O cliente fez uma análise inicial sobre o seguinte problema:
{analysis.get('problem_description', '')}

Resultados da análise anterior:
{context_text}

Agora o cliente quer aprofundar ou esclarecer pontos específicos. Responda de forma clara, objetiva e acionável.
Use Markdown para formatar sua resposta com listas, negrito e seções quando apropriado."""

        message = client.messages.create(
            model="claude-sonnet-4-20250514",
            max_tokens=2048,
            system=system_prompt,
            messages=[
                {"role": "user", "content": request.message}
            ]
        )
        
        response_text = message.content[0].text
        
        # Salva a conversa no histórico da análise
        if "chat_history" not in analysis:
            analysis["chat_history"] = []
        
        analysis["chat_history"].append({
            "role": "user",
            "content": request.message,
            "timestamp": datetime.utcnow().isoformat()
        })
        analysis["chat_history"].append({
            "role": "assistant", 
            "content": response_text,
            "timestamp": datetime.utcnow().isoformat()
        })
        
        # Retorna resposta com informações de uso
        new_used = used + 1
        return {
            "response": response_text,
            "usage": {
                "used": new_used,
                "limit": limit,
                "remaining": limit - new_used if limit != -1 else -1
            }
        }
        
    except Exception as e:
        print(f"Erro no refino: {str(e)}")
        return {"response": f"Desculpe, ocorreu um erro ao processar sua pergunta. Por favor, tente novamente."}


@app.get("/api/v1/async/analyses/{analysis_id}/export/{format}")
async def export_analysis(
    analysis_id: str,
    format: str,
    tenant: TenantContext = Depends(get_tenant_context)
):
    """Exporta análise para PDF, PPTX ou DOCX"""
    from fastapi.responses import Response
    import io
    
    analysis = analyses_store.get(analysis_id)
    
    if not analysis:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="Análise não encontrada"
        )
    
    if analysis.get("tenant_id") != tenant.tenant_id:
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="Acesso negado"
        )
    
    # Verifica permissão de exportação pelo plano
    can_export, error_msg = BillingService.can_export(tenant.tenant_id, format)
    if not can_export:
        raise HTTPException(
            status_code=status.HTTP_402_PAYMENT_REQUIRED,
            detail=error_msg
        )
    
    # Prepara conteúdo
    problem = analysis.get("problem_description", "")
    results = analysis.get("results", {})
    created_at = analysis.get("created_at", "")
    chat_history = analysis.get("chat_history", [])
    
    if format == "pdf":
        try:
            from reportlab.lib import colors
            from reportlab.lib.pagesizes import A4
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
            from reportlab.lib.units import inch
            
            buffer = io.BytesIO()
            doc = SimpleDocTemplate(buffer, pagesize=A4, topMargin=0.5*inch, bottomMargin=0.5*inch)
            styles = getSampleStyleSheet()
            
            # Custom styles
            title_style = ParagraphStyle('Title', parent=styles['Heading1'], fontSize=18, spaceAfter=20)
            heading_style = ParagraphStyle('Heading', parent=styles['Heading2'], fontSize=14, spaceAfter=10, textColor=colors.HexColor('#6366f1'))
            body_style = ParagraphStyle('Body', parent=styles['Normal'], fontSize=10, spaceAfter=8, leading=14)
            
            story = []
            
            # Title
            story.append(Paragraph("Relatório de Análise Estratégica", title_style))
            story.append(Paragraph(f"Gerado em: {created_at[:10] if created_at else 'N/A'}", body_style))
            story.append(Spacer(1, 20))
            
            # Problem
            story.append(Paragraph("Problema Analisado", heading_style))
            story.append(Paragraph(problem[:1000], body_style))
            story.append(Spacer(1, 15))
            
            # Results by agent
            agent_names = {
                'analyst': 'Analista de Negócios',
                'commercial': 'Especialista Comercial',
                'financial': 'Especialista Financeiro',
                'market': 'Especialista de Mercado',
                'reviewer': 'Diagnóstico Executivo'
            }
            
            for agent_key, agent_name in agent_names.items():
                content = results.get(agent_key, "")
                if content:
                    story.append(Paragraph(agent_name, heading_style))
                    # Limpa markdown básico para PDF
                    clean_content = content.replace('**', '').replace('##', '').replace('###', '').replace('- ', '• ')[:2000]
                    story.append(Paragraph(clean_content, body_style))
                    story.append(Spacer(1, 10))
            
            # Adiciona histórico de refinamento (chat)
            if chat_history:
                story.append(Spacer(1, 20))
                story.append(Paragraph("Refinamento da Análise", heading_style))
                for msg in chat_history:
                    role = "Pergunta" if msg.get("role") == "user" else "Resposta"
                    content = msg.get("content", "")[:1500]
                    clean_content = content.replace('**', '').replace('##', '').replace('###', '').replace('- ', '• ')
                    story.append(Paragraph(f"<b>{role}:</b>", body_style))
                    story.append(Paragraph(clean_content, body_style))
                    story.append(Spacer(1, 8))
            
            doc.build(story)
            buffer.seek(0)
            
            return Response(
                content=buffer.getvalue(),
                media_type="application/pdf",
                headers={"Content-Disposition": f"attachment; filename=analise-{analysis_id}.pdf"}
            )
        except ImportError:
            raise HTTPException(status_code=500, detail="Biblioteca reportlab não instalada")
    
    elif format == "docx":
        try:
            from docx import Document
            from docx.shared import Inches, Pt
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            
            doc = Document()
            
            # Title
            title = doc.add_heading('Relatório de Análise Estratégica', 0)
            title.alignment = WD_ALIGN_PARAGRAPH.CENTER
            
            doc.add_paragraph(f"Gerado em: {created_at[:10] if created_at else 'N/A'}")
            doc.add_paragraph()
            
            # Problem
            doc.add_heading('Problema Analisado', level=1)
            doc.add_paragraph(problem)
            
            # Results
            agent_names = {
                'analyst': 'Analista de Negócios',
                'commercial': 'Especialista Comercial',
                'financial': 'Especialista Financeiro',
                'market': 'Especialista de Mercado',
                'reviewer': 'Diagnóstico Executivo'
            }
            
            for agent_key, agent_name in agent_names.items():
                content = results.get(agent_key, "")
                if content:
                    doc.add_heading(agent_name, level=1)
                    doc.add_paragraph(content[:3000])
            
            # Adiciona histórico de refinamento (chat)
            if chat_history:
                doc.add_heading('Refinamento da Análise', level=1)
                for msg in chat_history:
                    role = "Pergunta" if msg.get("role") == "user" else "Resposta"
                    content = msg.get("content", "")[:2000]
                    p = doc.add_paragraph()
                    p.add_run(f"{role}: ").bold = True
                    p.add_run(content)
            
            buffer = io.BytesIO()
            doc.save(buffer)
            buffer.seek(0)
            
            return Response(
                content=buffer.getvalue(),
                media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                headers={"Content-Disposition": f"attachment; filename=analise-{analysis_id}.docx"}
            )
        except ImportError:
            raise HTTPException(status_code=500, detail="Biblioteca python-docx não instalada")
    
    elif format == "pptx":
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
            
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
            
            # Title slide
            slide_layout = prs.slide_layouts[6]  # Blank
            slide = prs.slides.add_slide(slide_layout)
            
            # Add title box
            left = Inches(0.5)
            top = Inches(2.5)
            width = Inches(12.333)
            height = Inches(2)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = "Análise Estratégica"
            p.font.size = Pt(44)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
            
            # Problem slide
            slide = prs.slides.add_slide(slide_layout)
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = "Problema Analisado"
            p.font.size = Pt(28)
            p.font.bold = True
            
            txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5))
            tf2 = txBox2.text_frame
            tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            p2.text = problem[:500]
            p2.font.size = Pt(16)
            
            # Agent slides
            agent_names = {
                'reviewer': 'Diagnóstico Executivo',
                'analyst': 'Análise de Negócios',
                'commercial': 'Estratégia Comercial',
                'financial': 'Análise Financeira',
                'market': 'Análise de Mercado',
            }
            
            for agent_key, agent_name in agent_names.items():
                content = results.get(agent_key, "")
                if content:
                    slide = prs.slides.add_slide(slide_layout)
                    
                    # Title
                    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1))
                    tf = txBox.text_frame
                    p = tf.paragraphs[0]
                    p.text = agent_name
                    p.font.size = Pt(28)
                    p.font.bold = True
                    
                    # Content
                    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5))
                    tf2 = txBox2.text_frame
                    tf2.word_wrap = True
                    p2 = tf2.paragraphs[0]
                    # Limpa e trunca
                    clean = content.replace('**', '').replace('##', '').replace('###', '')[:800]
                    p2.text = clean
                    p2.font.size = Pt(14)
            
            # Adiciona slides de refinamento (chat)
            if chat_history:
                slide = prs.slides.add_slide(slide_layout)
                txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1))
                tf = txBox.text_frame
                p = tf.paragraphs[0]
                p.text = "Refinamento da Análise"
                p.font.size = Pt(28)
                p.font.bold = True
                
                # Agrupa perguntas e respostas
                chat_text = ""
                for msg in chat_history[:6]:  # Limita a 3 pares (6 mensagens)
                    role = "P:" if msg.get("role") == "user" else "R:"
                    content = msg.get("content", "")[:300]
                    clean = content.replace('**', '').replace('##', '').replace('###', '')
                    chat_text += f"{role} {clean}\n\n"
                
                txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5))
                tf2 = txBox2.text_frame
                tf2.word_wrap = True
                p2 = tf2.paragraphs[0]
                p2.text = chat_text[:1500]
                p2.font.size = Pt(12)
            
            buffer = io.BytesIO()
            prs.save(buffer)
            buffer.seek(0)
            
            return Response(
                content=buffer.getvalue(),
                media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                headers={"Content-Disposition": f"attachment; filename=analise-{analysis_id}.pptx"}
            )
        except ImportError:
            raise HTTPException(status_code=500, detail="Biblioteca python-pptx não instalada")
    
    else:
        raise HTTPException(status_code=400, detail=f"Formato '{format}' não suportado. Use: pdf, docx, pptx")


@app.get("/api/v1/async/analyses/{analysis_id}")
async def get_analysis(
    analysis_id: str,
    tenant: TenantContext = Depends(get_tenant_context)
):
    """Retorna detalhe de uma análise"""
    analysis = analyses_store.get(analysis_id)
    
    if not analysis:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="Análise não encontrada"
        )
    
    if analysis.get("tenant_id") != tenant.tenant_id:
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="Acesso negado"
        )
    
    return analysis


# ============================================================================
# ENDPOINTS: BILLING
# ============================================================================

@app.get("/api/v1/billing/plans")
async def get_billing_plans():
    """Retorna planos disponíveis"""
    return {
        "plans": [
            {
                "tier": "free",
                "name": "Gratuito",
                "price": 0,
                "features": ["5 análises/mês", "3 agentes", "Suporte por email"],
                "limits": {"analyses_per_month": 5, "agents": 3}
            },
            {
                "tier": "pro",
                "name": "Profissional",
                "price": 99,
                "features": ["50 análises/mês", "5 agentes", "Suporte prioritário", "Exportação PDF"],
                "limits": {"analyses_per_month": 50, "agents": 5}
            },
            {
                "tier": "enterprise",
                "name": "Empresarial",
                "price": 299,
                "features": ["Análises ilimitadas", "5 agentes", "Suporte dedicado", "API access", "White-label"],
                "limits": {"analyses_per_month": -1, "agents": 5}
            }
        ],
        "currency": "BRL"
    }


@app.get("/api/v1/billing/status", response_model=BillingStatus)
async def get_billing_status(tenant: TenantContext = Depends(get_tenant_context)):
    """Retorna status de billing do tenant"""
    return BillingService.get_status(tenant.tenant_id)


# ============================================================================
# HEALTH CHECK
# ============================================================================

@app.get("/health")
async def health_check():
    """Health check para monitoramento"""
    return {
        "status": "healthy",
        "timestamp": datetime.utcnow().isoformat()
    }


if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)
