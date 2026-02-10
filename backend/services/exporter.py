"""
Exportador de análises para diferentes formatos.
"""

from typing import Dict, Any
from datetime import datetime
from io import BytesIO


class AnalysisExporter:
    """Exporta análises para Markdown, PDF e PowerPoint."""
    
    @staticmethod
    def to_markdown(data: Dict[str, Any]) -> str:
        """Exporta análise para Markdown."""
        problem = data.get("problem", "N/A")
        business_type = data.get("business_type", "N/A")
        timestamp = data.get("timestamp", datetime.now())
        results = data.get("results", {})
        
        if isinstance(timestamp, datetime):
            timestamp_str = timestamp.strftime("%d/%m/%Y %H:%M")
        else:
            timestamp_str = str(timestamp)
        
        md = f"""# Análise de Negócio

**Data:** {timestamp_str}  
**Tipo:** {business_type}

## Problema Analisado

{problem}

## Resultados da Análise

"""
        
        agent_names = {
            "analyst": "📊 Analista de Negócios",
            "commercial": "💼 Especialista Comercial",
            "financial": "💰 Especialista Financeiro",
            "market": "📈 Especialista de Mercado",
            "executive": "👔 Resumo Executivo",
        }
        
        for key, title in agent_names.items():
            if key in results:
                md += f"### {title}\n\n{results[key]}\n\n"
        
        return md
    
    @staticmethod
    def to_pdf(data: Dict[str, Any], filename: str) -> bytes:
        """Exporta análise para PDF."""
        try:
            from reportlab.lib.pagesizes import letter
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.lib.units import inch
            
            buffer = BytesIO()
            doc = SimpleDocTemplate(buffer, pagesize=letter)
            styles = getSampleStyleSheet()
            story = []
            
            # Título
            title_style = ParagraphStyle(
                'Title',
                parent=styles['Heading1'],
                fontSize=18,
                spaceAfter=12
            )
            story.append(Paragraph("Análise de Negócio", title_style))
            story.append(Spacer(1, 0.2 * inch))
            
            # Problema
            story.append(Paragraph("Problema Analisado", styles['Heading2']))
            story.append(Paragraph(data.get("problem", "N/A"), styles['Normal']))
            story.append(Spacer(1, 0.2 * inch))
            
            # Resultados
            results = data.get("results", {})
            agent_names = {
                "analyst": "Analista de Negócios",
                "commercial": "Especialista Comercial",
                "financial": "Especialista Financeiro",
                "market": "Especialista de Mercado",
                "executive": "Resumo Executivo",
            }
            
            for key, title in agent_names.items():
                if key in results and results[key]:
                    story.append(Paragraph(title, styles['Heading2']))
                    # Limitar tamanho para evitar problemas
                    content = str(results[key])[:5000]
                    story.append(Paragraph(content, styles['Normal']))
                    story.append(Spacer(1, 0.2 * inch))
            
            doc.build(story)
            return buffer.getvalue()
            
        except ImportError:
            # Fallback: retorna markdown como texto
            return AnalysisExporter.to_markdown(data).encode('utf-8')
    
    @staticmethod
    def to_ppt(data: Dict[str, Any], filename: str) -> bytes:
        """Exporta análise para PowerPoint."""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            
            prs = Presentation()
            
            # Slide de título
            title_slide = prs.slides.add_slide(prs.slide_layouts[0])
            title_slide.shapes.title.text = "Análise de Negócio"
            title_slide.placeholders[1].text = data.get("business_type", "")
            
            # Slide do problema
            problem_slide = prs.slides.add_slide(prs.slide_layouts[1])
            problem_slide.shapes.title.text = "Problema Analisado"
            problem_slide.placeholders[1].text = data.get("problem", "N/A")[:500]
            
            # Slides dos resultados
            results = data.get("results", {})
            agent_names = {
                "analyst": "Analista de Negócios",
                "commercial": "Especialista Comercial",
                "financial": "Especialista Financeiro",
                "market": "Especialista de Mercado",
                "executive": "Resumo Executivo",
            }
            
            for key, title in agent_names.items():
                if key in results and results[key]:
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    slide.shapes.title.text = title
                    # Limitar tamanho do texto
                    content = str(results[key])[:1000]
                    slide.placeholders[1].text = content
            
            buffer = BytesIO()
            prs.save(buffer)
            return buffer.getvalue()
            
        except ImportError:
            # Fallback: retorna markdown como texto
            return AnalysisExporter.to_markdown(data).encode('utf-8')
