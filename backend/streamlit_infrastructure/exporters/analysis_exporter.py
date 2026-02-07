"""Exportadores para análises de negócio."""

import io
from typing import Dict, Optional
from datetime import datetime


class AnalysisExporter:
    """Exporta análises em diferentes formatos."""
    
    @staticmethod
    def to_markdown(analysis_data: Dict) -> str:
        """
        Exporta análise como Markdown (one-pager).
        
        Args:
            analysis_data: Dicionário com dados da análise
        
        Returns:
            String com conteúdo Markdown
        """
        lines = []
        
        # Cabeçalho
        lines.append("# 🎯 Análise Estratégica de Negócio")
        lines.append("")
        lines.append(f"**Data**: {analysis_data.get('timestamp', datetime.now()).strftime('%d/%m/%Y %H:%M')}")
        lines.append(f"**Tipo de Negócio**: {analysis_data.get('business_type', 'N/A')}")
        lines.append(f"**Profundidade**: {analysis_data.get('analysis_depth', 'N/A')}")
        lines.append("")
        
        # Problema
        lines.append("## 📋 Problema/Oportunidade")
        lines.append("")
        lines.append(analysis_data.get('problem', 'N/A'))
        lines.append("")
        
        # Análises por agente
        results = analysis_data.get('results', {})
        
        if results.get('analyst'):
            lines.append("## 🔍 Análise de Negócio")
            lines.append("")
            lines.append(results['analyst'])
            lines.append("")
        
        if results.get('commercial'):
            lines.append("## 💼 Estratégia Comercial")
            lines.append("")
            lines.append(results['commercial'])
            lines.append("")
        
        if results.get('financial'):
            lines.append("## 💰 Análise Financeira")
            lines.append("")
            lines.append(results['financial'])
            lines.append("")
        
        if results.get('market'):
            lines.append("## 📊 Contexto de Mercado")
            lines.append("")
            lines.append(results['market'])
            lines.append("")
        
        if results.get('executive'):
            lines.append("## 👔 Decisão Executiva")
            lines.append("")
            lines.append(results['executive'])
            lines.append("")
        
        # Metadados
        metadata = results.get('metadata', {})
        if metadata:
            lines.append("---")
            lines.append("")
            lines.append("## 📊 Metadados da Execução")
            lines.append("")
            
            total_latency = sum(m.get('latency_ms', 0) for m in metadata.values())
            total_tokens = sum(m.get('tokens', 0) for m in metadata.values())
            total_cost = sum(m.get('cost_usd', 0) for m in metadata.values())
            
            lines.append(f"- **Latência Total**: {total_latency:.0f}ms")
            lines.append(f"- **Tokens Totais**: {total_tokens:,}")
            lines.append(f"- **Custo Total**: ${total_cost:.4f}")
            lines.append("")
            
            lines.append("### Por Agente:")
            for agent_name, agent_meta in metadata.items():
                lines.append(f"- **{agent_name}**: {agent_meta.get('latency_ms', 0):.0f}ms | {agent_meta.get('tokens', 0)} tokens | ${agent_meta.get('cost_usd', 0):.4f}")
        
        return "\n".join(lines)
    
    @staticmethod
    def _clean_text_for_pdf(text: str) -> str:
        """Limpa e formata texto para PDF."""
        if not text:
            return ""
        
        # Remove caracteres problemáticos para ReportLab
        text = text.replace('**', '')
        text = text.replace('###', '')
        text = text.replace('##', '')
        text = text.replace('#', '')
        text = text.replace('`', '')
        text = text.replace('*', '')
        
        # Converte quebras de linha para HTML
        paragraphs = text.split('\n\n')
        cleaned = []
        for p in paragraphs:
            p = p.strip()
            if p:
                # Converte linhas simples em quebras
                p = p.replace('\n', '<br/>')
                cleaned.append(p)
        
        return '<br/><br/>'.join(cleaned)
    
    @staticmethod
    def _extract_key_points(text: str, max_points: int = 5) -> list:
        """Extrai pontos-chave de um texto."""
        if not text:
            return []
        
        lines = text.split('\n')
        points = []
        
        for line in lines:
            line = line.strip()
            # Procura por linhas que parecem ser pontos importantes
            if line.startswith(('-', '•', '*', '1.', '2.', '3.', '4.', '5.')):
                point = line.lstrip('-•* 0123456789.')
                if len(point) > 10:
                    points.append(point.strip())
            elif ':' in line and len(line) < 200:
                points.append(line)
        
        return points[:max_points]
    
    @staticmethod
    def to_pdf(analysis_data: Dict, output_path: str) -> bytes:
        """
        Exporta análise como PDF profissional.
        
        Args:
            analysis_data: Dicionário com dados da análise
            output_path: Caminho para salvar PDF
        
        Returns:
            Bytes do PDF gerado
        """
        try:
            from reportlab.lib.pagesizes import A4
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.lib.units import inch, cm
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak, Table, TableStyle, ListFlowable, ListItem
            from reportlab.lib import colors
            from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_JUSTIFY
            
            # Cria documento em memória
            pdf_buffer = io.BytesIO()
            doc = SimpleDocTemplate(
                pdf_buffer,
                pagesize=A4,
                rightMargin=2*cm,
                leftMargin=2*cm,
                topMargin=2*cm,
                bottomMargin=2*cm
            )
            
            # Estilos personalizados
            styles = getSampleStyleSheet()
            
            title_style = ParagraphStyle(
                'CustomTitle',
                parent=styles['Heading1'],
                fontSize=22,
                textColor=colors.HexColor('#1a365d'),
                spaceAfter=20,
                spaceBefore=0,
                alignment=TA_CENTER,
                fontName='Helvetica-Bold'
            )
            
            section_style = ParagraphStyle(
                'SectionHeading',
                parent=styles['Heading2'],
                fontSize=14,
                textColor=colors.white,
                spaceAfter=10,
                spaceBefore=15,
                fontName='Helvetica-Bold',
                backColor=colors.HexColor('#2d3748'),
                borderPadding=8,
                leftIndent=0,
                rightIndent=0,
            )
            
            subsection_style = ParagraphStyle(
                'SubSection',
                parent=styles['Heading3'],
                fontSize=12,
                textColor=colors.HexColor('#2d3748'),
                spaceAfter=8,
                spaceBefore=12,
                fontName='Helvetica-Bold'
            )
            
            body_style = ParagraphStyle(
                'BodyText',
                parent=styles['Normal'],
                fontSize=10,
                textColor=colors.HexColor('#2d3748'),
                spaceAfter=8,
                spaceBefore=4,
                alignment=TA_JUSTIFY,
                leading=14,
                fontName='Helvetica'
            )
            
            highlight_style = ParagraphStyle(
                'Highlight',
                parent=styles['Normal'],
                fontSize=10,
                textColor=colors.HexColor('#1a365d'),
                spaceAfter=6,
                spaceBefore=4,
                backColor=colors.HexColor('#e2e8f0'),
                borderPadding=6,
                fontName='Helvetica'
            )
            
            meta_style = ParagraphStyle(
                'MetaInfo',
                parent=styles['Normal'],
                fontSize=9,
                textColor=colors.HexColor('#718096'),
                spaceAfter=4,
                fontName='Helvetica'
            )
            
            # Conteúdo
            story = []
            
            # === CABEÇALHO ===
            story.append(Paragraph("Análise Estratégica de Negócio", title_style))
            story.append(Spacer(1, 0.3*cm))
            
            # Metadados em tabela
            timestamp = analysis_data.get('timestamp', datetime.now())
            if isinstance(timestamp, str):
                timestamp_str = timestamp
            else:
                timestamp_str = timestamp.strftime('%d/%m/%Y às %H:%M')
            
            meta_data = [
                ['Data:', timestamp_str],
                ['Tipo de Negócio:', analysis_data.get('business_type', 'N/A')],
                ['Profundidade:', analysis_data.get('analysis_depth', 'Padrão')],
            ]
            
            meta_table = Table(meta_data, colWidths=[3*cm, 12*cm])
            meta_table.setStyle(TableStyle([
                ('FONTNAME', (0, 0), (0, -1), 'Helvetica-Bold'),
                ('FONTNAME', (1, 0), (1, -1), 'Helvetica'),
                ('FONTSIZE', (0, 0), (-1, -1), 9),
                ('TEXTCOLOR', (0, 0), (-1, -1), colors.HexColor('#4a5568')),
                ('VALIGN', (0, 0), (-1, -1), 'TOP'),
                ('BOTTOMPADDING', (0, 0), (-1, -1), 4),
            ]))
            story.append(meta_table)
            story.append(Spacer(1, 0.5*cm))
            
            # === PROBLEMA ===
            story.append(Paragraph("PROBLEMA / OPORTUNIDADE", section_style))
            problem_text = analysis_data.get('problem', 'N/A')
            # Remove dados anexados do problema se existir
            if '====' in problem_text:
                problem_text = problem_text.split('====')[0].strip()
            story.append(Paragraph(AnalysisExporter._clean_text_for_pdf(problem_text), highlight_style))
            story.append(Spacer(1, 0.3*cm))
            
            # === ANÁLISES ===
            results = analysis_data.get('results', {})
            
            sections = [
                ('executive', 'DIAGNÓSTICO EXECUTIVO', True),
                ('analyst', 'ANÁLISE DE NEGÓCIO', False),
                ('commercial', 'ESTRATÉGIA COMERCIAL', False),
                ('financial', 'ANÁLISE FINANCEIRA', False),
                ('market', 'CONTEXTO DE MERCADO', False),
            ]
            
            for key, title, is_main in sections:
                content = results.get(key, '')
                if content:
                    if is_main:
                        story.append(PageBreak())
                    
                    story.append(Paragraph(title, section_style))
                    
                    # Limpa e formata o conteúdo
                    clean_content = AnalysisExporter._clean_text_for_pdf(content)
                    
                    # Divide em parágrafos menores para melhor leitura
                    paragraphs = clean_content.split('<br/><br/>')
                    for para in paragraphs[:15]:  # Limita quantidade de parágrafos
                        if para.strip():
                            story.append(Paragraph(para.strip(), body_style))
                    
                    story.append(Spacer(1, 0.3*cm))
            
            # Constrói PDF
            doc.build(story)
            
            # Retorna bytes
            pdf_buffer.seek(0)
            return pdf_buffer.getvalue()
        
        except ImportError:
            raise ImportError("reportlab é necessário para exportar PDF. Instale com: pip install reportlab")
    
    @staticmethod
    def _clean_text_for_ppt(text: str, max_chars: int = 800) -> str:
        """Limpa texto para PowerPoint."""
        if not text:
            return ""
        
        # Remove formatação markdown
        text = text.replace('**', '')
        text = text.replace('###', '')
        text = text.replace('##', '')
        text = text.replace('#', '')
        text = text.replace('`', '')
        text = text.replace('*', '')
        
        # Limita tamanho
        if len(text) > max_chars:
            text = text[:max_chars] + "..."
        
        return text.strip()
    
    @staticmethod
    def _split_into_bullets(text: str, max_bullets: int = 6) -> list:
        """Divide texto em bullet points."""
        if not text:
            return []
        
        bullets = []
        lines = text.split('\n')
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
            
            # Remove prefixos de lista
            line = line.lstrip('-•* 0123456789.')
            line = line.strip()
            
            if len(line) > 20:  # Ignora linhas muito curtas
                # Limita tamanho de cada bullet
                if len(line) > 150:
                    line = line[:147] + "..."
                bullets.append(line)
        
        return bullets[:max_bullets]
    
    @staticmethod
    def to_ppt(analysis_data: Dict, output_path: str) -> bytes:
        """
        Exporta análise como PowerPoint profissional.
        
        Args:
            analysis_data: Dicionário com dados da análise
            output_path: Caminho para salvar PPTX
        
        Returns:
            Bytes do PPTX gerado
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt, Cm
            from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
            from pptx.dml.color import RGBColor
            
            # Cores do tema
            DARK_BLUE = RGBColor(26, 54, 93)
            LIGHT_GRAY = RGBColor(241, 245, 249)
            WHITE = RGBColor(255, 255, 255)
            DARK_GRAY = RGBColor(45, 55, 72)
            
            # Cria apresentação
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
            
            # === SLIDE 1: CAPA ===
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            
            # Background gradient (simulado com shape)
            background = slide.shapes.add_shape(
                1, Inches(0), Inches(0), prs.slide_width, prs.slide_height
            )
            background.fill.solid()
            background.fill.fore_color.rgb = DARK_BLUE
            background.line.fill.background()
            
            # Título principal
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = "Análise Estratégica de Negócio"
            p.font.size = Pt(44)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
            
            # Subtítulo
            subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
            tf = subtitle_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"Tipo: {analysis_data.get('business_type', 'N/A')}"
            p.font.size = Pt(20)
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
            
            # Data
            timestamp = analysis_data.get('timestamp', datetime.now())
            if isinstance(timestamp, str):
                date_str = timestamp
            else:
                date_str = timestamp.strftime('%d/%m/%Y')
            
            date_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12.333), Inches(0.5))
            tf = date_box.text_frame
            p = tf.paragraphs[0]
            p.text = date_str
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(200, 200, 200)
            p.alignment = PP_ALIGN.CENTER
            
            # === SLIDE 2: PROBLEMA ===
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            
            # Header bar
            header = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(1.2))
            header.fill.solid()
            header.fill.fore_color.rgb = DARK_BLUE
            header.line.fill.background()
            
            header_text = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12), Inches(0.6))
            tf = header_text.text_frame
            p = tf.paragraphs[0]
            p.text = "Problema / Oportunidade"
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = WHITE
            
            # Conteúdo do problema
            problem_text = analysis_data.get('problem', 'N/A')
            if '====' in problem_text:
                problem_text = problem_text.split('====')[0].strip()
            
            content_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.8), Inches(5))
            tf = content_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = AnalysisExporter._clean_text_for_ppt(problem_text, 600)
            p.font.size = Pt(16)
            p.font.color.rgb = DARK_GRAY
            p.line_spacing = 1.3
            
            # === SLIDES DE ANÁLISE ===
            results = analysis_data.get('results', {})
            
            sections = [
                ('executive', 'Diagnóstico Executivo'),
                ('analyst', 'Análise de Negócio'),
                ('commercial', 'Estratégia Comercial'),
                ('financial', 'Análise Financeira'),
                ('market', 'Contexto de Mercado'),
            ]
            
            for key, title in sections:
                content = results.get(key, '')
                if not content:
                    continue
                
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                
                # Header bar
                header = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(1.2))
                header.fill.solid()
                header.fill.fore_color.rgb = DARK_BLUE
                header.line.fill.background()
                
                header_text = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12), Inches(0.6))
                tf = header_text.text_frame
                p = tf.paragraphs[0]
                p.text = title
                p.font.size = Pt(28)
                p.font.bold = True
                p.font.color.rgb = WHITE
                
                # Conteúdo como bullets
                bullets = AnalysisExporter._split_into_bullets(content)
                
                if bullets:
                    content_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.6), Inches(11.8), Inches(5.5))
                    tf = content_box.text_frame
                    tf.word_wrap = True
                    
                    for i, bullet in enumerate(bullets):
                        if i == 0:
                            p = tf.paragraphs[0]
                        else:
                            p = tf.add_paragraph()
                        
                        p.text = f"• {bullet}"
                        p.font.size = Pt(14)
                        p.font.color.rgb = DARK_GRAY
                        p.line_spacing = 1.5
                        p.space_after = Pt(8)
                else:
                    # Se não conseguiu extrair bullets, usa texto limpo
                    content_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.6), Inches(11.8), Inches(5.5))
                    tf = content_box.text_frame
                    tf.word_wrap = True
                    p = tf.paragraphs[0]
                    p.text = AnalysisExporter._clean_text_for_ppt(content, 1000)
                    p.font.size = Pt(13)
                    p.font.color.rgb = DARK_GRAY
                    p.line_spacing = 1.4
            
            # === SLIDE FINAL ===
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            
            background = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
            background.fill.solid()
            background.fill.fore_color.rgb = DARK_BLUE
            background.line.fill.background()
            
            thanks_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12.333), Inches(1.5))
            tf = thanks_box.text_frame
            p = tf.paragraphs[0]
            p.text = "Obrigado"
            p.font.size = Pt(48)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
            
            subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12.333), Inches(1))
            tf = subtitle_box.text_frame
            p = tf.paragraphs[0]
            p.text = "Análise gerada pelo Consultor Executivo Multi-Agentes"
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(200, 200, 200)
            p.alignment = PP_ALIGN.CENTER
            
            # Salva em memória
            ppt_buffer = io.BytesIO()
            prs.save(ppt_buffer)
            ppt_buffer.seek(0)
            return ppt_buffer.getvalue()
        
        except ImportError:
            raise ImportError("python-pptx é necessário para exportar PowerPoint. Instale com: pip install python-pptx")
