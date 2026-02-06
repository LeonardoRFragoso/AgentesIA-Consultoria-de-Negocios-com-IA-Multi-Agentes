"""
Exportadores de artefatos executivos profissionais.
"""

import json
from typing import Optional
from datetime import datetime
from abc import ABC, abstractmethod

from core.executive_summary import ExecutiveSummary, ExecutiveReport, ExecutiveFormat
from infrastructure.logging import get_logger

logger = get_logger(__name__)


class ExecutiveExporter(ABC):
    """Classe base para exportadores executivos"""
    
    @abstractmethod
    def export(self, summary: ExecutiveSummary, output_path: Optional[str] = None) -> str:
        """
        Exporta resumo executivo.
        
        Args:
            summary: ExecutiveSummary a exportar
            output_path: Caminho para salvar (opcional)
        
        Returns:
            Conteúdo exportado como string
        """
        pass
    
    def _validate_summary(self, summary: ExecutiveSummary) -> None:
        """Valida resumo executivo"""
        if not summary.title:
            raise ValueError("Título é obrigatório")
        
        if not summary.key_decision:
            raise ValueError("Decisão é obrigatória")
        
        if not summary.rationale:
            raise ValueError("Rationale é obrigatória")
        
        if not (0.0 <= summary.confidence_score <= 1.0):
            raise ValueError("Confiança deve estar entre 0.0 e 1.0")
        
        if not summary.action_items:
            logger.warning("Nenhuma ação definida", summary_title=summary.title)


class OnePagerExporter(ExecutiveExporter):
    """Exporta one-pager em Markdown"""
    
    def export(self, summary: ExecutiveSummary, output_path: Optional[str] = None) -> str:
        """
        Exporta como one-pager em Markdown.
        
        Args:
            summary: ExecutiveSummary a exportar
            output_path: Caminho para salvar arquivo (opcional)
        
        Returns:
            Conteúdo Markdown
        """
        self._validate_summary(summary)
        
        content = summary.to_one_pager()
        
        if output_path:
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(content)
            
            logger.info(
                event="one_pager_exported",
                message=f"One-pager exported to {output_path}",
                summary_title=summary.title
            )
        
        return content


class PDFExporter(ExecutiveExporter):
    """Exporta relatório em PDF"""
    
    def export(self, report: ExecutiveReport, output_path: str) -> str:
        """
        Exporta como PDF formal.
        
        Args:
            report: ExecutiveReport a exportar
            output_path: Caminho para salvar PDF
        
        Returns:
            Mensagem de sucesso
        """
        self._validate_summary(report.summary)
        
        try:
            from reportlab.lib.pagesizes import letter, A4
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.lib.units import inch
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak, Table, TableStyle
            from reportlab.lib import colors
            
            # Cria documento
            doc = SimpleDocTemplate(
                output_path,
                pagesize=A4,
                rightMargin=0.75*inch,
                leftMargin=0.75*inch,
                topMargin=1*inch,
                bottomMargin=0.75*inch
            )
            
            # Estilos
            styles = getSampleStyleSheet()
            title_style = ParagraphStyle(
                'CustomTitle',
                parent=styles['Heading1'],
                fontSize=24,
                textColor=colors.HexColor('#1f4788'),
                spaceAfter=30,
                alignment=1  # Center
            )
            
            heading_style = ParagraphStyle(
                'CustomHeading',
                parent=styles['Heading2'],
                fontSize=14,
                textColor=colors.HexColor('#1f4788'),
                spaceAfter=12,
                spaceBefore=12
            )
            
            # Conteúdo
            story = []
            
            # Capa
            story.append(Paragraph(report.summary.title, title_style))
            story.append(Spacer(1, 0.2*inch))
            story.append(Paragraph(
                f"<b>Data:</b> {report.summary.date.strftime('%d/%m/%Y')}<br/>"
                f"<b>Responsável:</b> {report.summary.owner}",
                styles['Normal']
            ))
            story.append(Spacer(1, 0.5*inch))
            
            # Sumário executivo
            story.append(Paragraph("SUMÁRIO EXECUTIVO", heading_style))
            story.append(Paragraph(report.executive_summary_text, styles['Normal']))
            story.append(Spacer(1, 0.3*inch))
            
            # Contexto
            story.append(Paragraph("CONTEXTO", heading_style))
            story.append(Paragraph(report.background, styles['Normal']))
            story.append(Spacer(1, 0.3*inch))
            
            # Decisão
            story.append(Paragraph("DECISÃO", heading_style))
            story.append(Paragraph(
                f"<b>{report.summary.key_decision}</b><br/>"
                f"<i>Confiança: {report.summary.confidence_score:.0%}</i>",
                styles['Normal']
            ))
            story.append(Spacer(1, 0.3*inch))
            
            # Rationale
            story.append(Paragraph("RATIONALE", heading_style))
            story.append(Paragraph(report.summary.rationale, styles['Normal']))
            story.append(Spacer(1, 0.3*inch))
            
            # Ações
            if report.summary.action_items:
                story.append(Paragraph("AÇÕES IMEDIATAS", heading_style))
                for action in report.summary.action_items:
                    action_text = f"<b>{action.description}</b><br/>"
                    action_text += f"Responsável: {action.owner}"
                    if action.due_date:
                        action_text += f" | Prazo: {action.due_date}"
                    story.append(Paragraph(action_text, styles['Normal']))
                    story.append(Spacer(1, 0.1*inch))
                story.append(Spacer(1, 0.2*inch))
            
            # Riscos
            if report.summary.risks:
                story.append(Paragraph("RISCOS E MITIGAÇÃO", heading_style))
                for risk in report.summary.risks:
                    risk_text = f"<b>{risk.risk}</b><br/>"
                    risk_text += f"Probabilidade: {risk.probability} | Impacto: {risk.impact}"
                    if risk.mitigation:
                        risk_text += f"<br/>Mitigação: {risk.mitigation}"
                    story.append(Paragraph(risk_text, styles['Normal']))
                    story.append(Spacer(1, 0.1*inch))
                story.append(Spacer(1, 0.2*inch))
            
            # Próximos passos
            story.append(Paragraph("PRÓXIMOS PASSOS", heading_style))
            if report.summary.review_date:
                story.append(Paragraph(f"Revisar em {report.summary.review_date}", styles['Normal']))
            else:
                story.append(Paragraph("Monitorar implementação e coletar métricas", styles['Normal']))
            
            # Constrói PDF
            doc.build(story)
            
            logger.info(
                event="pdf_exported",
                message=f"PDF exported to {output_path}",
                summary_title=report.summary.title
            )
            
            return f"PDF exportado com sucesso: {output_path}"
        
        except ImportError:
            logger.error(
                event="pdf_export_failed",
                message="reportlab não está instalado",
                summary_title=report.summary.title,
                error="ImportError"
            )
            raise ImportError("reportlab é necessário para exportar PDF. Instale com: pip install reportlab")


class PPTExporter(ExecutiveExporter):
    """Exporta estrutura de apresentação em PowerPoint"""
    
    def export(self, report: ExecutiveReport, output_path: str) -> str:
        """
        Exporta estrutura de apresentação em PowerPoint.
        
        Args:
            report: ExecutiveReport a exportar
            output_path: Caminho para salvar PPTX
        
        Returns:
            Mensagem de sucesso
        """
        self._validate_summary(report.summary)
        
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
            from pptx.dml.color import RGBColor
            
            # Cria apresentação
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            
            # Obtém estrutura de slides
            slides_structure = report.get_ppt_structure()
            
            for slide_data in slides_structure:
                if slide_data["type"] == "cover":
                    self._add_cover_slide(prs, slide_data)
                elif slide_data["type"] == "highlight":
                    self._add_highlight_slide(prs, slide_data)
                elif slide_data["type"] == "list":
                    self._add_list_slide(prs, slide_data)
                else:  # text
                    self._add_text_slide(prs, slide_data)
            
            # Salva apresentação
            prs.save(output_path)
            
            logger.info(
                event="ppt_exported",
                message=f"PowerPoint exported to {output_path}",
                summary_title=report.summary.title,
                slide_count=len(slides_structure)
            )
            
            return f"PowerPoint exportado com sucesso: {output_path}"
        
        except ImportError:
            logger.error(
                event="ppt_export_failed",
                message="python-pptx não está instalado",
                summary_title=report.summary.title,
                error="ImportError"
            )
            raise ImportError("python-pptx é necessário para exportar PowerPoint. Instale com: pip install python-pptx")
    
    def _add_cover_slide(self, prs, slide_data):
        """Adiciona slide de capa"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Fundo
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(31, 71, 136)
        
        # Título
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(2))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        p = title_frame.paragraphs[0]
        p.text = slide_data["title"]
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(1))
        subtitle_frame = subtitle_box.text_frame
        p = subtitle_frame.paragraphs[0]
        p.text = slide_data["subtitle"]
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
    
    def _add_highlight_slide(self, prs, slide_data):
        """Adiciona slide com destaque"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        # Título
        title = slide.shapes.title
        title.text = slide_data["title"]
        
        # Conteúdo destacado
        content_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(3))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        p = content_frame.paragraphs[0]
        p.text = slide_data["content"]
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(31, 71, 136)
        p.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        if slide_data.get("subtitle"):
            subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(1))
            subtitle_frame = subtitle_box.text_frame
            p = subtitle_frame.paragraphs[0]
            p.text = slide_data["subtitle"]
            p.font.size = Pt(18)
            p.alignment = PP_ALIGN.CENTER
    
    def _add_text_slide(self, prs, slide_data):
        """Adiciona slide com texto"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        # Título
        title = slide.shapes.title
        title.text = slide_data["title"]
        
        # Conteúdo
        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5.5))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        p = content_frame.paragraphs[0]
        p.text = slide_data["content"]
        p.font.size = Pt(14)
        p.level = 0
    
    def _add_list_slide(self, prs, slide_data):
        """Adiciona slide com lista"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        # Título
        title = slide.shapes.title
        title.text = slide_data["title"]
        
        # Lista
        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5.5))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        for i, item in enumerate(slide_data["content"]):
            if i == 0:
                p = content_frame.paragraphs[0]
            else:
                p = content_frame.add_paragraph()
            
            p.text = f"• {item}"
            p.font.size = Pt(14)
            p.level = 0
            p.space_before = Pt(6)


class ExecutiveExporterFactory:
    """Factory para criar exportadores"""
    
    @staticmethod
    def create(format: ExecutiveFormat) -> ExecutiveExporter:
        """
        Cria exportador baseado no formato.
        
        Args:
            format: ExecutiveFormat desejado
        
        Returns:
            Exportador apropriado
        """
        if format == ExecutiveFormat.ONE_PAGER:
            return OnePagerExporter()
        elif format == ExecutiveFormat.PDF:
            return PDFExporter()
        elif format == ExecutiveFormat.PPT:
            return PPTExporter()
        else:
            raise ValueError(f"Formato desconhecido: {format}")
